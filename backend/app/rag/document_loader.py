"""
Document Loader — Extract text from PDF, DOCX, PPTX, and TXT files.
Handles various file formats and cleans extracted text.
"""
import os
import re
from typing import List, Dict, Any


def extract_text_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    """Extract text from PDF file, returning per-page content."""
    from PyPDF2 import PdfReader

    pages = []
    reader = PdfReader(file_path)
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({
                "content": text.strip(),
                "metadata": {"page_number": i + 1, "source_type": "PDF"}
            })
    return pages


def extract_text_from_docx(file_path: str) -> List[Dict[str, Any]]:
    """Extract text from DOCX file, returning per-paragraph content."""
    from docx import Document

    doc = Document(file_path)
    sections = []
    current_section = {"content": "", "metadata": {"section": "Main", "source_type": "DOCX"}}
    current_heading = "Main"

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Detect headings
        if para.style and para.style.name and "Heading" in para.style.name:
            # Save previous section if it has content
            if current_section["content"].strip():
                sections.append(current_section)
            current_heading = text
            current_section = {
                "content": "",
                "metadata": {"section": current_heading, "source_type": "DOCX"}
            }
        else:
            current_section["content"] += text + "\n"

    # Don't forget the last section
    if current_section["content"].strip():
        sections.append(current_section)

    return sections


def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """Extract text from PPTX file, returning per-slide content."""
    from pptx import Presentation

    slides = []
    prs = Presentation(file_path)

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            slides.append({
                "content": "\n".join(slide_text),
                "metadata": {"slide_number": i + 1, "source_type": "PPTX"}
            })

    return slides


def extract_text_from_txt(file_path: str) -> List[Dict[str, Any]]:
    """Extract text from plain text file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    return [{
        "content": content.strip(),
        "metadata": {"source_type": "TXT"}
    }]


def clean_text(text: str) -> str:
    """Clean and normalize extracted text."""
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    # Remove special characters that break LLM parsing
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
    # Normalize unicode
    text = text.strip()
    return text


def load_document(file_path: str) -> List[Dict[str, Any]]:
    """
    Main entry point — detect file type and extract text.
    Returns a list of content sections with metadata.
    """
    ext = os.path.splitext(file_path)[1].lower()

    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".txt": extract_text_from_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: {ext}. Supported: {list(extractors.keys())}")

    raw_sections = extractor(file_path)

    # Clean all extracted text
    for section in raw_sections:
        section["content"] = clean_text(section["content"])

    # Filter out empty sections
    sections = [s for s in raw_sections if len(s["content"]) > 20]

    if not sections:
        raise ValueError(f"No meaningful text extracted from {file_path}")

    return sections
